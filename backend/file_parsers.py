"""
Modular file parsers for extracting text from various file formats
Supports: DOCX, PPTX, XLSX, CSV, PDF, TXT
"""

import io
import re
from typing import Optional
from pathlib import Path

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False


class FileParser:
    """Base class for file parsers"""
    
    @staticmethod
    def extract_text(file_content: bytes, filename: str) -> Optional[str]:
        """
        Extract text from file content based on file extension
        
        Args:
            file_content: Raw bytes of the file
            filename: Original filename (for extension detection)
        
        Returns:
            Extracted text or None if parsing fails
        """
        file_ext = Path(filename).suffix.lower()
        
        parsers = {
            '.docx': FileParser._parse_docx,
            '.pptx': FileParser._parse_pptx,
            '.xlsx': FileParser._parse_xlsx,
            '.csv': FileParser._parse_csv,
            '.pdf': FileParser._parse_pdf,
            '.txt': FileParser._parse_text,
        }
        
        parser = parsers.get(file_ext)
        if not parser:
            print(f"Unsupported file type: {file_ext}")
            return None
        
        try:
            return parser(file_content)
        except Exception as e:
            print(f"Error parsing {filename}: {e}")
            return None
    
    @staticmethod
    def _parse_docx(file_content: bytes) -> str:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            print("python-docx not available")
            return ""
        
        doc = Document(io.BytesIO(file_content))
        text_parts = []
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        return "\n\n".join(text_parts)
    
    @staticmethod
    def _parse_pptx(file_content: bytes) -> str:
        """Extract text from PPTX file"""
        if not PPTX_AVAILABLE:
            print("python-pptx not available")
            return ""
        
        presentation = Presentation(io.BytesIO(file_content))
        text_parts = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_parts = [f"Slide {slide_num}:"]
            
            for shape in slide.shapes:
                # Extract text from text boxes
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text)
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_parts.append(" | ".join(row_text))
            
            if len(slide_parts) > 1:  # More than just the slide number
                text_parts.append("\n".join(slide_parts))
        
        return "\n\n---\n\n".join(text_parts)
    
    @staticmethod
    def _parse_xlsx(file_content: bytes) -> str:
        """Extract text from XLSX file"""
        if not PANDAS_AVAILABLE:
            print("pandas not available")
            return ""
        
        excel_file = io.BytesIO(file_content)
        text_parts = []
        
        try:
            # Read all sheets
            xls = pd.ExcelFile(excel_file)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name, engine='openpyxl')
                
                # Create descriptive text representation
                sheet_parts = [f"Sheet: {sheet_name}"]
                
                # Add column headers
                headers = [str(col) for col in df.columns if pd.notna(col)]
                if headers:
                    sheet_parts.append(f"Columns: {', '.join(headers)}")
                
                # Add data rows as text descriptions
                for idx, row in df.iterrows():
                    row_text = []
                    for col in df.columns:
                        value = row[col]
                        if pd.notna(value):
                            row_text.append(f"{col}: {value}")
                    
                    if row_text:
                        sheet_parts.append(f"Row {idx + 1}: " + " | ".join(row_text))
                
                text_parts.append("\n".join(sheet_parts))
        
        except Exception as e:
            print(f"Error parsing XLSX: {e}")
            return ""
        
        return "\n\n---\n\n".join(text_parts)
    
    @staticmethod
    def _parse_csv(file_content: bytes) -> str:
        """Extract text from CSV file"""
        if not PANDAS_AVAILABLE:
            print("pandas not available")
            return ""
        
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(io.BytesIO(file_content), encoding=encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                print("Could not decode CSV with any standard encoding")
                return ""
            
            text_parts = []
            
            # Add column headers
            headers = [str(col) for col in df.columns]
            text_parts.append(f"Columns: {', '.join(headers)}")
            
            # Add data rows
            for idx, row in df.iterrows():
                row_text = []
                for col in df.columns:
                    value = row[col]
                    if pd.notna(value):
                        row_text.append(f"{col}: {value}")
                
                if row_text:
                    text_parts.append(f"Row {idx + 1}: " + " | ".join(row_text))
            
            return "\n".join(text_parts)
        
        except Exception as e:
            print(f"Error parsing CSV: {e}")
            return ""
    
    @staticmethod
    def _parse_pdf(file_content: bytes) -> str:
        """Extract text from PDF file"""
        text_parts = []
        
        # Try PyMuPDF first (more robust)
        if PYMUPDF_AVAILABLE:
            try:
                pdf_doc = fitz.open(stream=file_content, filetype="pdf")
                for page_num in range(len(pdf_doc)):
                    page = pdf_doc[page_num]
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(f"Page {page_num + 1}:\n{text}")
                pdf_doc.close()
                return "\n\n---\n\n".join(text_parts)
            except Exception as e:
                print(f"PyMuPDF failed: {e}")
        
        # Fallback to PyPDF2
        if PYPDF2_AVAILABLE:
            try:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    if text.strip():
                        text_parts.append(f"Page {page_num}:\n{text}")
                return "\n\n---\n\n".join(text_parts)
            except Exception as e:
                print(f"PyPDF2 failed: {e}")
        
        print("No PDF parser available")
        return ""
    
    @staticmethod
    def _parse_text(file_content: bytes) -> str:
        """Extract text from plain text file"""
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252', 'ascii']
        
        for encoding in encodings:
            try:
                return file_content.decode(encoding)
            except UnicodeDecodeError:
                continue
        
        # Last resort: decode with errors ignored
        return file_content.decode('utf-8', errors='ignore')


def extract_text_from_bytes(file_content: bytes, filename: str) -> str:
    """
    Convenience function to extract text from file bytes
    
    Args:
        file_content: Raw file bytes
        filename: Original filename
    
    Returns:
        Extracted text or empty string
    """
    result = FileParser.extract_text(file_content, filename)
    return result if result else ""

